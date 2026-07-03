import os

# Cap CPU threads for the native math libraries before they load, so indexing
# can't peg every core and overheat the machine. Override with RAG_MAX_THREADS.
RAG_MAX_THREADS = max(1, int(os.getenv("RAG_MAX_THREADS", "4")))
for _var in ("OMP_NUM_THREADS", "OPENBLAS_NUM_THREADS", "MKL_NUM_THREADS",
             "VECLIB_MAXIMUM_THREADS", "NUMEXPR_NUM_THREADS"):
    os.environ.setdefault(_var, str(RAG_MAX_THREADS))
os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")

import re
import time
import httpx
import json
import numpy as np
import openpyxl
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from dotenv import load_dotenv
import chromadb
from chromadb.utils import embedding_functions

# Load .env file relative to this backend/ directory to ensure env variables
# are properly loaded regardless of the execution working directory.
_dotenv_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), ".env")
load_dotenv(_dotenv_path)

# Fallback to local 'resources' folder in the project root if RESOURCES_DIR is not set.
_default_resources_dir = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "resources"
)
RESOURCES_DIR = os.getenv("RESOURCES_DIR", _default_resources_dir)
DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY", "")
DEEPSEEK_API_BASE = os.getenv("DEEPSEEK_API_BASE", "https://api.deepseek.com/v1")
CHROMA_DB_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "chroma_db")

_ocr_engine = None
_onnx_threads_capped = False

def _cap_onnxruntime_threads():
    """onnxruntime defaults to one intra-op thread per core and ignores OMP_NUM_THREADS;
    neither chromadb nor rapidocr exposes the setting, so cap it at session creation."""
    global _onnx_threads_capped
    if _onnx_threads_capped:
        return
    _onnx_threads_capped = True
    try:
        import onnxruntime as ort
        orig_init = ort.InferenceSession.__init__

        def capped_init(self, *args, sess_options=None, **kwargs):
            so = sess_options or ort.SessionOptions()
            if so.intra_op_num_threads == 0:
                so.intra_op_num_threads = RAG_MAX_THREADS
                so.inter_op_num_threads = 1
            orig_init(self, *args, sess_options=so, **kwargs)

        ort.InferenceSession.__init__ = capped_init
    except Exception as e:
        print(f"Could not cap onnxruntime threads: {e}")

def _ocr_image_bytes(image_bytes):
    """OCR an embedded document image; returns extracted text or '' (never raises)."""
    global _ocr_engine
    try:
        if _ocr_engine is None:
            _cap_onnxruntime_threads()
            from rapidocr_onnxruntime import RapidOCR
            _ocr_engine = RapidOCR()
        import cv2
        arr = cv2.imdecode(np.frombuffer(image_bytes, np.uint8), cv2.IMREAD_COLOR)
        if arr is None:
            return ""
        # OCR cost grows with pixel area; diagram text stays legible at 2000px
        h, w = arr.shape[:2]
        if max(h, w) > 2000:
            scale = 2000 / max(h, w)
            arr = cv2.resize(arr, (int(w * scale), int(h * scale)), interpolation=cv2.INTER_AREA)
        result, _ = _ocr_engine(arr)
        lines = [r[1] for r in (result or [])]
        text = "\n".join(lines)
        return text if len(text.strip()) >= 10 else ""
    except Exception as e:
        print(f"OCR skipped for embedded image: {e}")
        return ""

class DocumentChunk:
    def __init__(self, chunk_id, filename, content, metadata=None):
        self.id = chunk_id
        self.filename = filename
        self.content = content
        self.metadata = metadata or {}

class ChromaSearchEngine:
    def __init__(self):
        self.chunks = []
        self.client = None
        self.collection = None
        self.embedding_function = None

    def initialize_chroma(self):
        if self.client is not None:
            return

        print(f"Initializing Persistent ChromaDB Client at: {CHROMA_DB_PATH}")
        _cap_onnxruntime_threads()
        self.client = chromadb.PersistentClient(path=CHROMA_DB_PATH)
        
        # Use default lightweight SentenceTransformer embeddings (downloads all-MiniLM-L6-v2)
        print("Loading default Chroma SentenceTransformer embedding function...")
        self.embedding_function = embedding_functions.DefaultEmbeddingFunction()
        
        # Get or create collection
        self.collection = self.client.get_or_create_collection(
            name="integration_repository",
            embedding_function=self.embedding_function
        )

    def load_and_index_documents(self):
        """Sync the persisted index with resources/: re-parse (and re-OCR/re-embed)
        only files that are new or changed since they were indexed, drop chunks of
        deleted files, and leave everything else untouched."""
        sync_start = time.time()
        self.initialize_chroma()
        self.chunks = []

        if not os.path.exists(RESOURCES_DIR):
            print(f"Resources directory not found: {RESOURCES_DIR}")
            return

        # Fingerprints of what's already indexed, keyed by filename
        indexed = {}
        for meta in self.collection.get(include=["metadatas"]).get("metadatas") or []:
            fn = (meta or {}).get("filename")
            if fn:
                indexed[fn] = (meta.get("file_mtime"), meta.get("file_size"))

        on_disk = {}
        for filename in os.listdir(RESOURCES_DIR):
            filepath = os.path.join(RESOURCES_DIR, filename)
            if os.path.isdir(filepath) or filename.startswith('.'):
                continue
            if os.path.splitext(filename)[1].lower() not in ('.xlsx', '.docx', '.pptx', '.pdf'):
                continue
            st = os.stat(filepath)
            on_disk[filename] = (st.st_mtime, st.st_size)

        removed = [fn for fn in indexed if fn not in on_disk]
        for fn in removed:
            self.collection.delete(where={"filename": fn})

        stale = [fn for fn in on_disk if indexed.get(fn) != on_disk[fn]]
        reindexed = 0
        for filename in stale:
            filepath = os.path.join(RESOURCES_DIR, filename)
            ext = os.path.splitext(filename)[1].lower()
            before = len(self.chunks)
            try:
                if ext == '.xlsx':
                    self._parse_xlsx(filename, filepath, before)
                elif ext == '.docx':
                    self._parse_docx(filename, filepath, before)
                elif ext == '.pptx':
                    self._parse_pptx(filename, filepath, before)
                elif ext == '.pdf':
                    self._parse_pdf(filename, filepath, before)
            except Exception as e:
                # Keep the previously indexed version searchable
                print(f"Error parsing {filename}: {e}")
                del self.chunks[before:]
                continue

            new_chunks = self.chunks[before:]
            mtime, size = on_disk[filename]
            for j, chunk in enumerate(new_chunks):
                chunk.id = f"{filename}::{j}"
                chunk.metadata["filename"] = filename
                chunk.metadata["file_mtime"] = mtime
                chunk.metadata["file_size"] = size

            self.collection.delete(where={"filename": filename})
            batch_size = 100
            for i in range(0, len(new_chunks), batch_size):
                batch = new_chunks[i:i + batch_size]
                self.collection.add(
                    documents=[c.content for c in batch],
                    metadatas=[c.metadata for c in batch],
                    ids=[c.id for c in batch]
                )
            reindexed += 1

        # Hydrate the in-memory chunk list from the now-current collection
        self.chunks = []
        data = self.collection.get(include=["documents", "metadatas"])
        for cid, doc, meta in zip(data["ids"], data["documents"], data["metadatas"]):
            meta = meta or {}
            self.chunks.append(DocumentChunk(cid, meta.get("filename", "unknown"), doc, meta))

        print(f"Index sync: {len(on_disk) - len(stale)} unchanged, {reindexed} reindexed, "
              f"{len(removed)} removed; {len(self.chunks)} chunks total "
              f"({time.time() - sync_start:.1f}s)")

    def _parse_xlsx(self, filename, filepath, start_idx):
        wb = openpyxl.load_workbook(filepath, read_only=True)
        idx = start_idx
        for sheetname in wb.sheetnames:
            sheet = wb[sheetname]
            rows = []
            for r_idx, row in enumerate(sheet.iter_rows(values_only=True)):
                if any(c is not None for c in row):
                    row_str = " | ".join(str(c) if c is not None else "" for c in row)
                    rows.append(f"Row {r_idx+1}: {row_str}")
            
            chunk_size = 15
            for i in range(0, len(rows), chunk_size):
                chunk_rows = rows[i:i+chunk_size]
                content = f"Document: {filename}\nSheet: {sheetname}\n" + "\n".join(chunk_rows)
                metadata = {
                    "sheet": sheetname,
                    "type": "spreadsheet",
                    "rows": f"{i+1}-{i+len(chunk_rows)}"
                }
                self.chunks.append(DocumentChunk(f"chunk_{idx}", filename, content, metadata))
                idx += 1

    def _parse_docx(self, filename, filepath, start_idx):
        doc = Document(filepath)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        
        for t_idx, table in enumerate(doc.tables):
            table_rows = []
            for r_idx, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                table_rows.append(" | ".join(cells))
            paragraphs.append(f"[Table {t_idx+1} Content]:\n" + "\n".join(table_rows))
            
        idx = start_idx
        current_chunk = []
        current_len = 0
        
        for p in paragraphs:
            current_chunk.append(p)
            current_len += len(p)
            if current_len >= 1200:
                content = f"Document: {filename}\n" + "\n\n".join(current_chunk)
                self.chunks.append(DocumentChunk(f"chunk_{idx}", filename, content, {"type": "document"}))
                idx += 1
                current_chunk = current_chunk[-1:]
                current_len = len(current_chunk[0])
                
        if current_chunk:
            content = f"Document: {filename}\n" + "\n\n".join(current_chunk)
            self.chunks.append(DocumentChunk(f"chunk_{idx}", filename, content, {"type": "document"}))
            idx += 1

        # OCR embedded images (architecture diagrams etc.) so their content is retrievable.
        # Captions are paired by order: figure captions appear in the same sequence as images.
        from docx.parts.image import ImagePart
        captions = [p for p in paragraphs if p.lower().startswith(("figure", "diagram"))]
        images = sorted(
            (p for p in doc.part.related_parts.values() if isinstance(p, ImagePart)),
            key=lambda p: str(p.partname)
        )
        for i_idx, image in enumerate(images):
            ocr_text = _ocr_image_bytes(image.blob)
            if not ocr_text:
                continue
            caption = captions[i_idx] if i_idx < len(captions) else ""
            header = f"Reference diagram{' — ' + caption if caption else ''}"
            content = f"Document: {filename}\n{header}\nText extracted from the diagram image:\n{ocr_text}"
            self.chunks.append(DocumentChunk(f"chunk_{idx}", filename, content, {"type": "diagram"}))
            idx += 1

    def _parse_pptx(self, filename, filepath, start_idx):
        prs = Presentation(filepath)
        idx = start_idx
        for s_idx, slide in enumerate(prs.slides):
            slide_text = []
            title = f"Slide {s_idx+1}"
            if slide.shapes.title:
                title = slide.shapes.title.text
                slide_text.append(f"Title: {title}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if slide.shapes.title and shape.text == slide.shapes.title.text:
                        continue
                    slide_text.append(shape.text.strip())
            
            content = f"Document: {filename}\nSlide {s_idx+1}: {title}\n" + "\n".join(slide_text)
            self.chunks.append(DocumentChunk(f"chunk_{idx}", filename, content, {
                "slide": s_idx+1,
                "title": title,
                "type": "presentation"
            }))
            idx += 1

            # OCR pictures on the slide (diagrams, screenshots)
            for shape in slide.shapes:
                try:
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        ocr_text = _ocr_image_bytes(shape.image.blob)
                        if ocr_text:
                            content = (f"Document: {filename}\nReference diagram — slide {s_idx+1}: {title}\n"
                                       f"Text extracted from the diagram image:\n{ocr_text}")
                            self.chunks.append(DocumentChunk(f"chunk_{idx}", filename, content, {
                                "slide": s_idx+1,
                                "type": "diagram"
                            }))
                            idx += 1
                except Exception as e:
                    print(f"Skipping image on slide {s_idx+1} of {filename}: {e}")

    def _pdf_page_jpegs(self, page):
        jpegs = []
        try:
            xobjs = page["/Resources"]["/XObject"]
        except Exception:
            return jpegs
        for name in list(xobjs.keys()):
            try:
                obj = xobjs[name].get_object()
                if obj.get("/Subtype") != "/Image":
                    continue
                filters = obj.get("/Filter")
                if not isinstance(filters, list):
                    filters = [filters]
                if "/DCTDecode" not in [str(f) for f in filters]:
                    continue
                if obj.get("/Width", 0) < 200 or obj.get("/Height", 0) < 150:
                    continue
                jpegs.append(obj.get_data())
            except Exception:
                continue
        return jpegs

    def _parse_pdf(self, filename, filepath, start_idx):
        reader = PdfReader(filepath)
        idx = start_idx
        for p_idx, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                content = f"Document: {filename}\nPage {p_idx+1}:\n{text}"
                self.chunks.append(DocumentChunk(f"chunk_{idx}", filename, content, {
                    "page": p_idx+1,
                    "type": "pdf"
                }))
                idx += 1

            # OCR embedded images (architecture diagrams rendered as pictures).
            # Only large JPEG (DCTDecode) XObjects: their raw stream is directly decodable,
            # while pypdf's decoding of exotic formats (CCITT fax etc.) can segfault.
            for image_bytes in self._pdf_page_jpegs(page):
                try:
                    ocr_text = _ocr_image_bytes(image_bytes)
                    if ocr_text:
                        content = (f"Document: {filename}\nReference diagram — page {p_idx+1}\n"
                                   f"Text extracted from the diagram image:\n{ocr_text}")
                        self.chunks.append(DocumentChunk(f"chunk_{idx}", filename, content, {
                            "page": p_idx+1,
                            "type": "diagram"
                        }))
                        idx += 1
                except Exception as e:
                    print(f"Skipping image on page {p_idx+1} of {filename}: {e}")

    def search(self, query, top_k=5, where=None):
        self.initialize_chroma()
        if not self.collection:
            return []

        try:
            results = self.collection.query(
                query_texts=[query],
                n_results=top_k,
                where=where
            )
            
            retrieved_chunks = []
            if results and 'documents' in results and len(results['documents']) > 0:
                docs = results['documents'][0]
                metadatas = results['metadatas'][0]
                ids = results['ids'][0]
                
                for idx in range(len(docs)):
                    meta = metadatas[idx]
                    filename = meta.get("filename", "unknown")
                    # Recreate DocumentChunk object
                    retrieved_chunks.append(DocumentChunk(
                        chunk_id=ids[idx],
                        filename=filename,
                        content=docs[idx],
                        metadata=meta
                    ))
            return retrieved_chunks
        except Exception as e:
            print(f"Error querying ChromaDB: {e}")
            return []

engine = ChromaSearchEngine()

def get_engine(force_reload=False):
    global engine
    if force_reload or not engine.chunks:
        engine.load_and_index_documents()
    return engine

SVG_RULES = """Rules for generating the SVG diagram:
MOST IMPORTANT — using reference diagrams from the context (marked 'Reference diagram' / 'Text extracted from the diagram image'):
- If a reference diagram covers the SAME scenario and systems as the request, mirror it faithfully: its tiers/layers, its exact API and system names, and its flow order.
- If a reference diagram is a generic pattern/template for the requested diagram style (e.g. an API-led connectivity template built around a different business domain), mirror its STRUCTURE ONLY: the tier stacking and ordering, one tier per horizontal band, the tier-to-tier flow direction, and backend systems at the bottom — but populate it with the APIs and systems from the user's scenario.
- Only fall back to generic API-led conventions where the context is silent.

1. Output exactly ONE self-contained, responsive `<svg>` element inside an xml code block (```xml ... ```). Do not include any HTML wrapping, raw text, or explanations inside the block.
2. SVG Design guidelines:
   - Use a modern dark theme background: '#0b101c' or transparent.
   - Use rounded rectangles (rx="8", ry="8") for API/system boxes (standard size: width="200" height="60").
   - Layer layout and y-positions:
     * Consumer Experience / Front-end (y=30): e.g. Mobile App, Web Browser
     * Experience APIs Layer (y=150): fill='#0e7490' (Cyan), stroke='#22d3ee', stroke-width='2'
     * Process APIs Layer (y=290): fill='#1d4ed8' (Blue), stroke='#60a5fa', stroke-width='2'
     * System APIs Layer (y=430): fill='#6d28d9' (Purple), stroke='#a78bfa', stroke-width='2'
     * Backend Systems Layer (y=570): fill='#334155' (Slate), stroke='#94a3b8', stroke-width='2'
   - Align nodes horizontally within each tier so they are neat. Space multiple nodes in a tier horizontally (e.g. x=100, x=350, x=600 for a width of 900).
   - Draw connector lines with arrowheads (using SVG `<line>` or `<path>` with `marker-end`) connecting the nodes between layers.
   - Place small, legible text labels (fill='#94a3b8', font-size='12px') next to lines indicating protocol/payload (e.g. HTTPS/JSON, JMS).
   - Text inside boxes: Use '<text text-anchor="middle" dominant-baseline="middle" fill="#ffffff" font-size="13px" font-family="system-ui, sans-serif">' and place them in the center of the box. Use multiple text lines or `<tspan>` if needed.
3. Keep the layout responsive using viewBox="0 0 900 660" with width="100%" and height="100%". Include standard markers for arrowheads at the top of the SVG:
   <defs>
     <marker id="arrow" viewBox="0 0 10 10" refX="6" refY="5" markerWidth="6" markerHeight="6" orient="auto-start-reverse">
       <path d="M 0 2 L 8 5 L 0 8 z" fill="#7dd3fc" />
     </marker>
   </defs>
"""

def compute_grounding(response_text: str, context_chunks: list):
    """
    Estimate what share of a response is grounded in the retrieved document chunks
    versus generated from the model's own knowledge.

    Each response segment (sentence or code block) is embedded with the same local
    SentenceTransformer used for indexing; its max cosine similarity against the
    retrieved chunks is mapped onto a 0-1 grounded score, then averaged weighted
    by segment length. This is a heuristic, not an exact attribution.
    """
    try:
        if not response_text or not context_chunks:
            return None
        engine.initialize_chroma()

        segments = []
        for piece in re.split(r'(```[\s\S]*?```)', response_text):
            if piece.startswith('```'):
                body = piece.strip('`').strip()
                if len(body) > 40:
                    segments.append(body[:1000])
            else:
                for sent in re.split(r'(?<=[.!?:])\s+|\n+', piece):
                    sent = re.sub(r'[#*>`]', ' ', sent).strip()
                    if len(sent) > 30:
                        segments.append(sent[:1000])
        if not segments:
            return None

        chunk_texts = [c.content[:2000] for c in context_chunks]
        seg_emb = np.array(engine.embedding_function(segments))
        ctx_emb = np.array(engine.embedding_function(chunk_texts))
        seg_emb = seg_emb / np.linalg.norm(seg_emb, axis=1, keepdims=True)
        ctx_emb = ctx_emb / np.linalg.norm(ctx_emb, axis=1, keepdims=True)
        max_sims = (seg_emb @ ctx_emb.T).max(axis=1)

        # Below LOW cosine similarity a segment counts as pure model knowledge,
        # above HIGH as fully document-grounded; linear in between.
        # Thresholds calibrated against this corpus: unrelated text ~0.05,
        # paraphrased doc content 0.3-0.6, near-verbatim 0.6+.
        LOW, HIGH = 0.20, 0.60
        emb_scores = np.clip((max_sims - LOW) / (HIGH - LOW), 0.0, 1.0)

        # Embeddings under-credit verbatim copies of long chunks (a single copied
        # sentence embeds far from the whole chunk), so also measure word-trigram
        # containment and take the stronger signal per segment.
        def shingles(text):
            words = re.findall(r'[a-z0-9]+', text.lower())
            return set(zip(words, words[1:], words[2:]))

        ctx_shingles = set()
        for c in context_chunks:
            ctx_shingles |= shingles(c.content)

        lex_scores = []
        for s in segments:
            seg_sh = shingles(s)
            lex_scores.append(len(seg_sh & ctx_shingles) / len(seg_sh) if seg_sh else 0.0)

        scores = np.maximum(emb_scores, np.array(lex_scores))
        weights = np.array([min(len(s), 400) for s in segments], dtype=float)
        rag_pct = int(round(float((scores * weights).sum() / weights.sum()) * 100))
        return {
            "rag_pct": rag_pct,
            "model_pct": 100 - rag_pct,
            "segments_scored": len(segments)
        }
    except Exception as e:
        print(f"Error computing grounding score: {e}")
        return None

async def ask_deepseek_llm(query: str, context_chunks: list, user_api_key: str = None, history: list = None) -> str:
    api_key = user_api_key or DEEPSEEK_API_KEY
    if not api_key:
        return "DeepSeek API Key is missing. Please provide your API Key either in the UI settings or backend config."

    context_str = "\n\n---\n\n".join([chunk.content for chunk in context_chunks])
    
    system_prompt = (
        "You are 'Integration Architect AI', an expert integration architect. "
        "Your task is to answer technical questions, design APIs, review compliance, and debug error logs "
        "using only the provided context from the enterprise integration documentation repository.\n\n"
        "Rules:\n"
        "1. Ground your answers strictly in the provided document context. If the document context doesn't contain "
        "the details, state that clearly but provide a reasonable response based on general integration best practices if helpful.\n"
        "2. Provide concrete code snippets (e.g. MuleSoft XML, DataWeave code, JSON schema, or REST endpoints) when requested.\n"
        "3. Maintain a highly professional, expert tone suitable for an Integration Architect.\n"
        "4. Do not mention document chunk IDs or metadata variables unless specifically relevant; refer to files by their actual names.\n"
        "5. When the user asks for an architecture, flow, or sequence diagram, respond with a valid, self-contained SVG diagram "
        "inside an xml code block (```xml ... ```). Draw a beautiful, professional, structured diagram representing the requested flow. "
        "For API-led connectivity, stack the layers top-to-bottom: Experience Layer (Top), Process Layer (Middle), "
        "System Layer (Bottom), and Backend/External Systems (Bottom-most). "
        "Use rounded rectangles (rx=\"8\") with custom colors to distinguish tiers:\n"
        "   - Experience APIs: fill '#0e7490' (Cyan), stroke '#22d3ee', stroke-width='2'\n"
        "   - Process APIs: fill '#1d4ed8' (Blue), stroke '#60a5fa', stroke-width='2'\n"
        "   - System APIs: fill '#6d28d9' (Purple), stroke '#a78bfa', stroke-width='2'\n"
        "   - Backend/External Systems: fill '#334155' (Slate), stroke '#94a3b8', stroke-width='2'\n"
        "Draw clear connector lines with arrowheads (using marker-end=\"url(#arrow)\") and protocol labels (like HTTPS/JSON or JMS). "
        "If the context contains a reference diagram ('Text extracted from the diagram image') for the SAME scenario, mirror its tiers and API names faithfully. "
        "Ensure the SVG includes viewBox=\"0 0 900 660\", width=\"100%\", height=\"100%\", and standard arrow marker defs.\n\n"
        "CONTEXT FROM ENTERPRISE INTEGRATION REPOSITORY:\n"
        f"{context_str}"
    )

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }

    messages = [{"role": "system", "content": system_prompt}]
    for m in (history or []):
        role = "assistant" if m.get("role") in ("assistant", "agent") else "user"
        content = m.get("content", "")
        if content:
            messages.append({"role": role, "content": content})
    messages.append({"role": "user", "content": query})

    payload = {
        "model": "deepseek-chat",
        "messages": messages,
        "temperature": 0.2,
        "max_tokens": 2048
    }

    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            response = await client.post(
                f"{DEEPSEEK_API_BASE}/chat/completions",
                headers=headers,
                json=payload
            )
            if response.status_code == 200:
                result = response.json()
                return result["choices"][0]["message"]["content"]
            else:
                return f"Error from DeepSeek API (Status Code {response.status_code}): {response.text}"
    except Exception as e:
        return f"Exception occurred while calling DeepSeek API: {str(e)}"
